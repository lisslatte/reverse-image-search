import os
import numpy as np
import cv2
from skimage.metrics import structural_similarity as ssim
from sklearn.decomposition import PCA
from sklearn.neighbors import NearestNeighbors
from tkinter import Tk, Toplevel, filedialog, Button, Label, Frame, Text, Scrollbar, ttk, messagebox
from PIL import Image, ImageTk, ExifTags
import docx
import PyPDF2
from pptx import Presentation
import io
import fitz  # PyMuPDF for PDF processing
import hashlib
import logging
from cryptography.fernet import Fernet
from docx import Document
import atexit
from sympy import content
from tkinter import font as tkfont

#Install required libraries
def install_packages():
    import subprocess
    import sys

    required_packages = [
        "numpy", "opencv-python", "scikit-image", "scikit-learn", "pillow", 
        "python-docx", "PyPDF2", "python-pptx", "pymupdf", "cryptography", "sympy"
    ]

    for package in required_packages:
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", package])
            print(f"Successfully installed {package}")
        except subprocess.CalledProcessError as e:
            print(f"Failed to install {package}. Error: {e}")

# Configure logging
logging.basicConfig(filename='app.log', level=logging.ERROR)

# Encryption key 
key = Fernet.generate_key()
cipher_suite = Fernet(key)

def encrypt_file(file_path):
    with open(file_path, 'rb') as file:
        file_data = file.read()
    encrypted_data = cipher_suite.encrypt(file_data)
    with open(file_path, 'wb') as file:
        file.write(encrypted_data)

def decrypt_file(file_path):
    with open(file_path, 'rb') as file:
        encrypted_data = file.read()
    decrypted_data = cipher_suite.decrypt(encrypted_data)
    with open(file_path, 'wb') as file:
        file.write(decrypted_data)

def load_images_from_folder(folder):
    if not os.path.isdir(folder):
        raise ValueError("The specified folder does not exist.")
    
    images = []
    image_paths = []
    valid_extensions = ('.jpg', '.jpeg', '.png')  # Supported image formats
    for filename in os.listdir(folder):
        if filename.lower().endswith(valid_extensions):
            img_path = os.path.join(folder, filename)
            if os.path.isfile(img_path):  # Ensure it's a file
                img = cv2.imread(img_path)
                if img is not None:
                    images.append(img)
                    image_paths.append(img_path)
                else:
                    logging.error(f"Failed to load image: {img_path}")
    return images, image_paths

def extract_features(images):
    features = []
    for img in images:
        if img is not None:
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            resized = cv2.resize(gray, (128, 128)).flatten()
            features.append(resized)
        else:
            logging.error("Encountered an empty image while extracting features.")
    return np.array(features)


def index_images(features):
    pca = PCA(n_components=50)
    features_pca = pca.fit_transform(features)
    nbrs = NearestNeighbors(n_neighbors=5, algorithm='auto').fit(features_pca)
    return nbrs, pca

def calculate_similarity(image1, image2):
    gray1 = cv2.cvtColor(image1, cv2.COLOR_BGR2GRAY)
    gray2 = cv2.cvtColor(image2, cv2.COLOR_BGR2GRAY)
    resized1 = cv2.resize(gray1, (128, 128))  # Ensure same dimensions
    resized2 = cv2.resize(gray2, (128, 128))  # Ensure same dimensions
    return ssim(resized1, resized2)

def calculate_color_histogram_similarity(image1, image2):
    # Calculate color histograms for each image
    hist1 = cv2.calcHist([image1], [0, 1, 2], None, [8, 8, 8], [0, 256, 0, 256, 0, 256])
    hist2 = cv2.calcHist([image2], [0, 1, 2], None, [8, 8, 8], [0, 256, 0, 256, 0, 256])
    
    # Normalize the histograms
    hist1 = cv2.normalize(hist1, hist1).flatten()
    hist2 = cv2.normalize(hist2, hist2).flatten()
    
    # Calculate the correlation between the histograms
    similarity = cv2.compareHist(hist1, hist2, cv2.HISTCMP_CORREL)
    return similarity

def get_image_metadata(image_path):
    try:
        with Image.open(image_path) as img:
            metadata = {
                'format': img.format,
                'mode': img.mode,
                'size': img.size,
                'filename': os.path.basename(img.filename),
                'info': img.info
            }
            if hasattr(img, '_getexif'):
                exifdata = img._getexif()
                if exifdata:
                    for tag, value in exifdata.items():
                        decoded = ExifTags.TAGS.get(tag, tag)
                        metadata[decoded] = value
        return metadata
    except Exception as e:
        logging.error(f"Failed to get metadata for image {image_path}: {e}")
        return {'error': str(e)}

def search_similar_images(image_path, nbrs, pca, image_paths, number_of_images=5):
    img = cv2.imread(image_path)
    if img is None:
        logging.error(f"Failed to load query image: {image_path}")
        raise ValueError(f"Could not load the query image: {image_path}")
    
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    resized = cv2.resize(gray, (128, 128)).flatten()
    query_feature = pca.transform([resized])
    distances, indices = nbrs.kneighbors(query_feature, n_neighbors=number_of_images)
    
    similar_images = []
    
    # Include the original query image
    query_metadata = get_image_metadata(image_path)
    similar_images.append((image_path, 100.0, query_metadata))
    
    for idx in indices[0]:
        similar_img_path = image_paths[idx]
        similar_img = cv2.imread(similar_img_path)
        if similar_img is not None:
            structural_similarity = calculate_similarity(img, similar_img)
            color_histogram_similarity = calculate_color_histogram_similarity(img, similar_img)
            combined_similarity = (structural_similarity + color_histogram_similarity) / 2
            metadata = get_image_metadata(similar_img_path)
            similar_images.append((similar_img_path, combined_similarity * 100, metadata))  # Convert to percentage
        else:
            logging.error(f"Failed to load similar image: {similar_img_path}")
    
    # Sort similar images by similarity in descending order
    similar_images.sort(key=lambda x: x[1], reverse=True)
    
    return similar_images


def extract_images_from_pdf_document(document_path):
    images = []
    try:
        doc = fitz.open(document_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image = Image.open(io.BytesIO(image_bytes))
                
                # Convert RGBA to RGB if necessary
                if image.mode == 'RGBA':
                    image = image.convert('RGB')
                
                img_path = f"temp_image_{page_num}_{img_index}.jpg"
                image.save(img_path)
                images.append(img_path)
    except Exception as e:
        logging.error(f"Failed to extract images from PDF: {e}")
        messagebox.showerror("Error", f"Failed to extract images from PDF: {e}")
    return images

def extract_images_from_word_document(document_path):
    images = []
    try:
        doc = Document(document_path)
        for rel in doc.part.rels.values():
            if 'image' in rel.target_ref:
                image_stream = rel.target_part.blob
                img = Image.open(io.BytesIO(image_stream))
                
                # Determine image file extension
                if img.format == 'JPEG':
                    file_ext = 'jpg'
                elif img.format == 'PNG':
                    file_ext = 'png'
                elif img.format == 'GIF':
                    file_ext = 'gif'
                else:
                    logging.warning(f"Unsupported image format: {img.format}")
                    continue
                
                # Construct image path
                img_name = os.path.basename(rel.target_ref)
                img_path = os.path.join(os.path.dirname(document_path), f"{img_name}.{file_ext}")
                
                # Save the image with original format
                img.save(img_path)
                
                images.append(img_path)
    except Exception as e:
        logging.error(f"Failed to extract images from Word document: {e}")
        messagebox.showerror("Error", f"Failed to extract images from Word document: {e}")
    return images

def extract_images_from_pptx_document(document_path):
    images = []
    try:
        prs = Presentation(document_path)
        for slide_num, slide in enumerate(prs.slides):
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):  # Check if the shape has an image
                    image = shape.image
                    image_bytes = image.blob  # Access the binary data of the image
                    file_ext = image.ext  # Get the file extension
                    
                    img_path = os.path.join(os.path.dirname(document_path), f"slide_{slide_num+1}_image_{shape_num+1}.{file_ext}")
                    
                    # Save the image with its original format
                    with open(img_path, 'wb') as f:
                        f.write(image_bytes)
                    
                    images.append(img_path)
    except Exception as e:
        logging.error(f"Failed to extract images from PowerPoint document: {e}")
        messagebox.showerror("Error", f"Failed to extract images from PowerPoint document: {e}")
    return images

def extract_images_from_document(document_path):
    images = []
    if document_path.endswith('.docx'):
        images = extract_images_from_word_document(document_path)
    elif document_path.endswith('.pdf'):
        images = extract_images_from_pdf_document(document_path)
    elif document_path.endswith('.pptx'):
        images = extract_images_from_pptx_document(document_path)
    else:
        raise ValueError("Unsupported document format.")
    
    if images:
        messagebox.showinfo("Images Found", f"Found {len(images)} images in the document.")
    else:
        messagebox.showwarning("No Images Found", "No images were found in the document.")
    
    return images

class ImageTab:
    def __init__(self, tab_control, img_path, nbrs, pca, image_paths):
        self.tab_control = tab_control
        self.tab = ttk.Frame(tab_control)
        tab_control.add(self.tab, text=os.path.basename(img_path))
        self.image_path = img_path
        self.similar_images = search_similar_images(img_path, nbrs, pca, image_paths, number_of_images=5)
        self.current_index = 0

        self.display_area = Frame(self.tab)
        self.display_area.pack()

        self.nav_frame = Frame(self.tab)
        self.nav_frame.pack(pady=10)
        
        self.prev_button = Button(self.nav_frame, text="Previous", command=self.show_previous)
        self.prev_button.pack(side='left', padx=10)

        self.next_button = Button(self.nav_frame, text="Next", command=self.show_next)
        self.next_button.pack(side='left', padx=10)

        self.clear_button = Button(self.nav_frame, text="Clear", command=self.clear_tab)
        self.clear_button.pack(side='left', padx=10)
        
        self.show_image()

    def show_image(self):
        for widget in self.display_area.winfo_children():
            widget.destroy()

        img_path, similarity, metadata = self.similar_images[self.current_index]

        img = Image.open(img_path)
        img.thumbnail((400, 400), Image.ANTIALIAS)
        img_tk = ImageTk.PhotoImage(img)

        panel = Label(self.display_area, image=img_tk)
        panel.image = img_tk  # Keep a reference to avoid garbage collection
        panel.pack()

        title = Label(self.display_area, text=f"Similarity: {similarity:.2f}%\n{metadata.get('filename', 'Unknown')}\n{metadata.get('size', 'Unknown')} {metadata.get('format', 'Unknown')}")
        title.pack()

        metadata_text = Text(self.display_area, height=10, width=50, wrap='word')
        metadata_text.insert('1.0', f"Detailed Metadata:\n\n")
        for key, value in metadata.items():
            metadata_text.insert('end', f"{key}: {value}\n")
        metadata_text.pack()

    def show_previous(self):
        if self.current_index > 0:
            self.current_index -= 1
            self.show_image()

    def show_next(self):
        if self.current_index < len(self.similar_images) - 1:
            self.current_index += 1
            self.show_image()

    def clear_tab(self):
        self.tab_control.forget(self.tab)

def hash_file(file_path):
    sha256_hash = hashlib.sha256()
    with open(file_path, "rb") as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()

def upload_and_search():
    global tab_control, nbrs, pca, image_paths, extracted_images
    file_path = filedialog.askopenfilename(filetypes=[("Image files", "*.jpg *.jpeg *.png"), ("Document files", "*.docx *.pdf *.pptx")])
    if file_path:
        try:
            if file_path.endswith(('.jpg', '.jpeg', '.png')):
                tab_control.add(ImageTab(tab_control, file_path, nbrs, pca, image_paths).tab)
            else:
                images = extract_images_from_document(file_path)
                if images:
                    extracted_images.extend(images)
                    for img_path in images:
                        tab_control.add(ImageTab(tab_control, img_path, nbrs, pca, image_paths).tab)
        except ValueError as e:
            logging.error(f"Unsupported file type: {e}")
            messagebox.showerror("Unsupported File Type", str(e))
        except Exception as e:
            logging.error(f"Failed to process the file: {e}")
            messagebox.showerror("Error", f"Failed to process the file: {e}")

def cleanup_extracted_images():
    global extracted_images
    for img_path in extracted_images:
        try:
            os.remove(img_path)
        except Exception as e:
            logging.error(f"Failed to delete extracted image {img_path}: {e}")

def show_text_file(filename, title):
        try:
            with open(filename, 'r') as file:
                content = file.read()
            show_text_window(content, title)
        except FileNotFoundError:
            messagebox.showerror("Error", f"{title} file not found.")

def show_text_window(content, title):
    window = Toplevel()
    window.title(title)
    text_widget = Text(window, wrap='word', height=30, width=80)
    text_widget.insert('1.0', content)
    text_widget.config(state='disabled')
    text_widget.pack(expand=1, fill='both')
    scrollbar = Scrollbar(window, command=text_widget.yview)
    scrollbar.pack(side='right', fill='y')
    text_widget.config(yscrollcommand=scrollbar.set)

def show_privacy_policy():
        show_text_file('privacy_policy.txt', 'Privacy Policy')

def show_tutorial():
        show_text_file('tutorial.txt', 'Tutorial')

def show_terms_conditions():
        show_text_file('terms_conditions.txt', 'Terms & Conditions')

def open_terms_and_conditions(event=None):
    show_terms_conditions()

def open_privacy_policy(event=None):
    show_privacy_policy()

def main():
    global tab_control, nbrs, pca, image_paths, extracted_images

    folder = 'images'  
    if not os.path.isdir(folder):
        logging.error("The specified folder does not exist.")
        messagebox.showerror("Error", "The specified folder does not exist.")
        return

    images, image_paths = load_images_from_folder(folder)
    features = extract_features(images)
    nbrs, pca = index_images(features)
    extracted_images = []

    root = Tk()
    root.title("I-Check (Image Similarity Search)")

    # Create a custom font
    title_font = tkfont.Font(family="Helvetica", size=24, weight="bold")   

    Label(root, text="I-CHECK", font=title_font).pack(pady=20)
    Label(root, text="Upload an image (JPG or PNG) or documents (PDF, WORD, PPTX) to find similar images").pack(pady=10)
    Button(root, text="Upload", command=upload_and_search).pack(pady=10)
    Button(root, text="Tutorial", command=show_tutorial).pack(pady=10)

    # Create a frame to hold the clickable labels
    bottom_frame = Frame(root)
    bottom_frame.pack(pady=10)
    
    # Terms and Conditions label
    terms_label = Label(bottom_frame, text="Terms and Conditions", fg="blue", cursor="hand2")
    terms_label.pack(side='left', padx=10)
    terms_label.bind("<Button-1>", open_terms_and_conditions)
    
    # Privacy Policy label
    privacy_label = Label(bottom_frame, text="Privacy Policy", fg="blue", cursor="hand2")
    privacy_label.pack(side='left', padx=10)
    privacy_label.bind("<Button-1>", open_privacy_policy)

    tab_control = ttk.Notebook(root)
    tab_control.pack(expand=1, fill='both')

    # Register cleanup function to delete extracted images on program exit
    atexit.register(cleanup_extracted_images)

    root.mainloop()

if __name__ == "__main__":
    install_packages()
    main()